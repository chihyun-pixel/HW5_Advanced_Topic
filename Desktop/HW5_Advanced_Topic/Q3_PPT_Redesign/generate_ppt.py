from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ========= Style A : Business / Minimal =========
prs = Presentation("original.pptx")

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(28)
                    run.font.name = "Arial"

prs.save("style_A.pptx")

# ========= Style B : Dark Theme =========
prs = Presentation("original.pptx")

for slide in prs.slides:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(30, 30, 30)

    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(26)
                    run.font.name = "Calibri"
                    run.font.color.rgb = RGBColor(255, 255, 255)

prs.save("style_B.pptx")
